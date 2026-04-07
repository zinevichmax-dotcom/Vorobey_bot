"""
Умный нормализатор PPTX v2.
Определяет «родной» стиль презентации и сбрасывает только аномалии.

Ключевое: учитывает тему (None = наследование от темы).
Если 90%+ runs наследуют стиль — любой явный шрифт с малой частотой = аномалия.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from collections import Counter
import json


def normalize_pptx(input_path: str, output_path: str) -> dict:
    """
    Нормализует PPTX: исправляет аномалии, сохраняет авторский стиль.
    """
    prs = Presentation(input_path)

    # Проход 1: статистика
    stats = _collect_stats(prs)

    # Определяем норму
    norm = _define_norm(stats)

    # Проход 2: фиксим
    fixes = _fix_anomalies(prs, norm)

    prs.save(output_path)

    return {
        "stats": {
            "slides": len(prs.slides),
            "total_runs": stats["total_runs"],
            "inherited_font_runs": stats["inherited_fonts"],
            "explicit_font_runs": stats["explicit_fonts"],
            "inherited_size_runs": stats["inherited_sizes"],
            "explicit_size_runs": stats["explicit_sizes"],
            "fonts_found": dict(stats["fonts"]),
            "sizes_found_pt": dict(stats["sizes_pt"]),
        },
        "norm": norm,
        "fixes": fixes,
        "total_fixes": len(fixes),
    }


def _collect_stats(prs) -> dict:
    """Проход 1: частотность + соотношение inherited/explicit."""
    fonts = Counter()
    sizes_pt = Counter()
    colors = Counter()
    total_runs = 0
    inherited_fonts = 0
    explicit_fonts = 0
    inherited_sizes = 0
    explicit_sizes = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            frames = []
            if shape.has_text_frame:
                frames.append(shape.text_frame)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        frames.append(cell.text_frame)

            for tf in frames:
                for para in tf.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        total_runs += 1

                        # Шрифты
                        if run.font.name:
                            explicit_fonts += 1
                            fonts[run.font.name] += 1
                        else:
                            inherited_fonts += 1

                        # Размеры
                        if run.font.size:
                            explicit_sizes += 1
                            pt = round(run.font.size / 12700, 1)
                            sizes_pt[pt] += 1
                        else:
                            inherited_sizes += 1

                        # Цвета
                        try:
                            if run.font.color and run.font.color.rgb:
                                colors[str(run.font.color.rgb)] += 1
                        except Exception:
                            pass

    return {
        "fonts": fonts,
        "sizes_pt": sizes_pt,
        "colors": colors,
        "total_runs": total_runs,
        "inherited_fonts": inherited_fonts,
        "explicit_fonts": explicit_fonts,
        "inherited_sizes": inherited_sizes,
        "explicit_sizes": explicit_sizes,
    }


def _define_norm(stats: dict) -> dict:
    """
    Определяет норму. Два режима:

    A) Тема доминирует (>70% inherited) → любой явный шрифт с count < 3 = аномалия
    B) Явные шрифты доминируют → классика: топ шрифты = норма, редкие = аномалия
    """
    fonts = stats["fonts"]
    sizes_pt = stats["sizes_pt"]
    total = stats["total_runs"]
    inherited_fonts = stats["inherited_fonts"]
    explicit_fonts = stats["explicit_fonts"]
    inherited_sizes = stats["inherited_sizes"]
    explicit_sizes = stats["explicit_sizes"]

    # --- Режим шрифтов ---
    if total > 0 and inherited_fonts / total > 0.7:
        font_mode = "theme_dominant"
        # Тема доминирует: родные = шрифты с count >= 5 ИЛИ >= 10% от explicit
        threshold = max(3, explicit_fonts * 0.1) if explicit_fonts > 0 else 3
        native_fonts = [f for f, c in fonts.items() if c >= threshold]
    else:
        font_mode = "explicit_dominant"
        # Явные доминируют: топ-5 + всё что >= 3
        native_fonts = list(
            set([f for f, _ in fonts.most_common(5)] + [f for f, c in fonts.items() if c >= 3])
        )

    # --- Режим размеров ---
    if total > 0 and inherited_sizes / total > 0.7:
        size_mode = "theme_dominant"
        # Тема доминирует: любой явный размер с count < 2 подозрителен
        # Но нужен контекст — если размер в разумном диапазоне, оставляем
        common_sizes = [s for s, c in sizes_pt.items() if c >= 2]
        if common_sizes:
            size_range = (min(common_sizes) * 0.7, max(common_sizes) * 1.3)
        else:
            # Нет common sizes — используем абсолютный диапазон
            size_range = (8.0, 60.0)
        # Одиночные размеры вне диапазона = аномалия
    else:
        size_mode = "explicit_dominant"
        common_sizes = [s for s, c in sizes_pt.items() if c >= 2]
        if not common_sizes and sizes_pt:
            common_sizes = list(sizes_pt.keys())
        if common_sizes:
            margin = max((max(common_sizes) - min(common_sizes)) * 0.3, 4)
            size_range = (max(min(common_sizes) - margin, 6), max(common_sizes) + margin)
        else:
            size_range = (8.0, 60.0)

    return {
        "font_mode": font_mode,
        "native_fonts": native_fonts,
        "size_mode": size_mode,
        "size_range_pt": [round(size_range[0], 1), round(size_range[1], 1)],
        "font_anomaly_rule": (
            "explicit font with count < threshold → reset to theme"
            if font_mode == "theme_dominant"
            else "font not in native list → reset to theme"
        ),
        "size_anomaly_rule": (f"explicit size outside {size_range[0]:.0f}-{size_range[1]:.0f}pt → reset"),
    }


def _fix_anomalies(prs, norm: dict) -> list:
    """Проход 2: исправляет аномалии."""
    fixes = []
    native_fonts = set(norm["native_fonts"])
    min_pt, max_pt = norm["size_range_pt"]
    font_mode = norm["font_mode"]

    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            frames = []
            if shape.has_text_frame:
                frames.append(("shape", shape.name, shape.text_frame))
            if shape.has_table:
                for ri, row in enumerate(shape.table.rows):
                    for ci, cell in enumerate(row.cells):
                        frames.append(("table", f"{shape.name}[{ri}][{ci}]", cell.text_frame))

            for source, name, tf in frames:
                for para in tf.paragraphs:
                    for run in para.runs:
                        text = run.text.strip()
                        if not text:
                            continue
                        preview = text[:40]

                        # --- Шрифт ---
                        if run.font.name:
                            is_alien = False
                            if font_mode == "theme_dominant":
                                # В теме: любой явный шрифт не в native = аномалия
                                is_alien = run.font.name not in native_fonts
                            else:
                                # Явные: шрифт не в native = аномалия
                                is_alien = native_fonts and run.font.name not in native_fonts

                            if is_alien:
                                fixes.append(
                                    {
                                        "slide": i + 1,
                                        "element": name,
                                        "type": "alien_font",
                                        "was": run.font.name,
                                        "action": "reset to theme",
                                        "text": preview,
                                    }
                                )
                                run.font.name = None

                        # --- Размер ---
                        if run.font.size:
                            pt = round(run.font.size / 12700, 1)
                            if pt < min_pt or pt > max_pt:
                                fixes.append(
                                    {
                                        "slide": i + 1,
                                        "element": name,
                                        "type": "abnormal_size",
                                        "was_pt": pt,
                                        "allowed_range": [min_pt, max_pt],
                                        "action": "reset to theme",
                                        "text": preview,
                                    }
                                )
                                run.font.size = None

    return fixes


# === CLI ===
if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Использование: python pptx_normalizer.py файл.pptx [output.pptx]")
        sys.exit(1)

    input_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else input_file.replace(".pptx", "_normalized.pptx")

    report = normalize_pptx(input_file, output_file)
    print(json.dumps(report, indent=2, default=str, ensure_ascii=False))
    print(f"\nСохранено: {output_file}")
