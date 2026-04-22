"""
Сборщик слайдов.
Входы:
  - style ("formal"/"corporate"/"bold")
  - layout_kind ("bullets"/"title"/...)
  - content dict с реальными значениями для плейсхолдеров

Алгоритм:
  1. Открыть master.pptx выбранного стиля
  2. Найти эталонный слайд нужного layout_kind (по индексу в master'е)
  3. Скопировать все shapes этого слайда в целевую презу
  4. В скопированных shapes заменить плейсхолдеры на реальный текст
  5. Если нужна картинка — заменить тестовое изображение на реальное
"""
import copy
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Inches
from pptx.oxml.ns import qn
from lxml import etree
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    from .layout_catalog import LAYOUT_KINDS
except ImportError:
    from layout_catalog import LAYOUT_KINDS


# Путь к папке masters/. По умолчанию — рядом с этим файлом.
# Можно переопределить через env REDESIGNER_MASTERS_DIR.
MASTERS_DIR = os.environ.get(
    "REDESIGNER_MASTERS_DIR",
    os.path.join(os.path.dirname(os.path.abspath(__file__)), "masters"),
)
MASTER_FILES = {
    "formal": "master_style_1_formal.pptx",
    "corporate": "master_style_2_corporate.pptx",
    "bold": "master_style_3_bold.pptx",
}


def _copy_shape_xml(src_shape, target_slide):
    """Копировать shape через XML в новый слайд."""
    src_elem = src_shape._element
    new_elem = copy.deepcopy(src_elem)
    target_slide.shapes._spTree.append(new_elem)
    return new_elem


def _clone_slide_with_rels(src_slide, target_prs, blank_layout):
    """
    Клонировать слайд целиком (shapes + media + rels) из source_prs в target_prs.
    Используется для _original — когда нужно сохранить слайд как есть.
    """
    target_slide = target_prs.slides.add_slide(blank_layout)

    # Background
    try:
        _copy_background(src_slide, target_slide)
    except Exception:
        pass

    # Копируем shapes и переносим связанные media (picture rels)
    src_part = src_slide.part
    target_part = target_slide.part

    # Map: old_rId -> new_rId
    rid_map = {}

    for shape in src_slide.shapes:
        src_elem = shape._element
        new_elem = copy.deepcopy(src_elem)

        # Найти все ссылки вида r:embed, r:link внутри скопированного элемента
        for attr_name in ("r:embed", "r:link"):
            # lxml ищет по qn()
            for elem in new_elem.iter():
                for a_key, a_val in list(elem.attrib.items()):
                    if a_key.endswith("}embed") or a_key.endswith("}link"):
                        old_rid = a_val
                        if old_rid in rid_map:
                            elem.set(a_key, rid_map[old_rid])
                            continue
                        # Взять target этого rel в source part
                        try:
                            related = src_part.related_part(old_rid)
                        except KeyError:
                            continue
                        # Добавить такое же media в target part
                        new_rid = target_part.relate_to(
                            related, src_part.rels[old_rid].reltype
                        )
                        rid_map[old_rid] = new_rid
                        elem.set(a_key, new_rid)

        target_slide.shapes._spTree.append(new_elem)

    return target_slide


def _copy_background(src_slide, target_slide):
    """Копировать background."""
    src_bg = src_slide.background._element
    target_bg = target_slide.background._element
    # Replace bg children
    for child in list(target_bg):
        target_bg.remove(child)
    for child in src_bg:
        target_bg.append(copy.deepcopy(child))


def _extract_placeholder_name(text: str) -> str | None:
    """Если текст выглядит как {NAME}, вернуть NAME."""
    text = text.strip()
    if text.startswith("{") and text.endswith("}"):
        return text[1:-1]
    return None


def _find_placeholders_in_xml(slide_elem) -> list[tuple]:
    """
    Найти все текстовые run'ы с плейсхолдерами вида {NAME}.
    Возвращает [(placeholder_name, run_xml_element), ...]
    """
    found = []
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    
    for run in slide_elem.iter(qn("a:r")):
        t_elem = run.find(qn("a:t"))
        if t_elem is not None and t_elem.text:
            name = _extract_placeholder_name(t_elem.text)
            if name:
                found.append((name, run))
    
    # Также паттерны внутри runs (которые не полностью плейсхолдер, например часть)
    for t_elem in slide_elem.iter(qn("a:t")):
        if t_elem.text and "{" in t_elem.text and "}" in t_elem.text:
            # Может быть частичная замена — оставим для find_and_replace
            pass
    
    return found


def _replace_in_text_elements(slide_elem, replacements: dict):
    """
    Пройти по всем a:t элементам и заменить {PLACEHOLDER} → значение.
    Работает со строчным поиском (robust против частичных плейсхолдеров).
    """
    for t_elem in slide_elem.iter(qn("a:t")):
        if t_elem.text:
            new_text = t_elem.text
            for key, value in replacements.items():
                placeholder = "{" + key + "}"
                if placeholder in new_text:
                    new_text = new_text.replace(placeholder, str(value))
            t_elem.text = new_text


def _replace_image_in_slide(slide, image_path: str) -> bool:
    """
    Заменить самую большую картинку на слайде (по площади) на image_path.
    Сохраняет aspect ratio исходной картинки, центрует в bounds placeholder'а.
    """
    if not image_path or not os.path.exists(image_path):
        return False
    
    largest = None
    largest_area = 0
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            area = (shape.width or 0) * (shape.height or 0)
            if area > largest_area:
                largest = shape
                largest_area = area
    
    if largest is None:
        return False
    
    ph_left, ph_top = largest.left, largest.top
    ph_width, ph_height = largest.width, largest.height
    
    # Вычислить размеры с сохранением aspect ratio
    from PIL import Image as PILImage
    try:
        with PILImage.open(image_path) as img:
            img_w, img_h = img.size
    except Exception:
        # fallback: без PIL — просто растягиваем как раньше
        sp = largest._element
        sp.getparent().remove(sp)
        slide.shapes.add_picture(image_path, ph_left, ph_top, width=ph_width, height=ph_height)
        return True
    
    # fit inside placeholder, keep ratio
    scale = min(ph_width / img_w, ph_height / img_h)
    new_w = int(img_w * scale)
    new_h = int(img_h * scale)
    # центровка
    new_left = ph_left + (ph_width - new_w) // 2
    new_top = ph_top + (ph_height - new_h) // 2
    
    sp = largest._element
    sp.getparent().remove(sp)
    slide.shapes.add_picture(image_path, new_left, new_top, width=new_w, height=new_h)
    return True


def _clear_missing_placeholders(slide_elem):
    """
    Убрать оставшиеся плейсхолдеры типа {XXX} (если не были заменены).
    """
    for t_elem in slide_elem.iter(qn("a:t")):
        if t_elem.text and "{" in t_elem.text and "}" in t_elem.text:
            # Удалить части вида {XXX}
            import re
            t_elem.text = re.sub(r'\{[A-Z_0-9]+\}', '', t_elem.text).strip()


def build_presentation(
    slides_spec: list[dict],
    style: str,
    output_path: str,
    source_pptx_path: str | None = None,
) -> str:
    """
    Собрать новую презентацию.
    slides_spec: список из {
        "layout_kind": str,
        "content": dict,
        "source_slide_idx": int (только для _original),
    }
    style: "formal" / "corporate" / "bold"
    output_path: куда сохранить
    source_pptx_path: нужен только если есть _original слайды для копирования
    """
    # Открыть master стиля
    master_path = os.path.join(MASTERS_DIR, MASTER_FILES[style])
    master_prs = Presentation(master_path)
    
    # Создать пустую презу с размерами master'а
    output_prs = Presentation()
    output_prs.slide_width = master_prs.slide_width
    output_prs.slide_height = master_prs.slide_height
    blank_layout = output_prs.slide_layouts[6]
    
    # Исходная преза (для _original слайдов)
    source_prs = Presentation(source_pptx_path) if source_pptx_path else None
    
    for spec in slides_spec:
        kind = spec["layout_kind"]
        content = spec.get("content", {})
        
        if kind == "_original":
            if source_prs and "source_slide_idx" in spec:
                src_slide = source_prs.slides[spec["source_slide_idx"]]
                # Правильное клонирование: копируем slide через XML и rels
                _clone_slide_with_rels(src_slide, output_prs, blank_layout)
            continue
        
        if kind not in LAYOUT_KINDS:
            # Неизвестный лейаут — пропустить или fallback
            continue
        
        # Найти эталонный слайд в master'е по индексу
        layout_idx = LAYOUT_KINDS.index(kind)
        if layout_idx >= len(master_prs.slides):
            continue
        
        master_slide = master_prs.slides[layout_idx]
        
        # Создать новый слайд в output
        target_slide = output_prs.slides.add_slide(blank_layout)
        
        # Скопировать background из master
        _copy_background(master_slide, target_slide)
        
        # Скопировать все shapes из master
        for shape in master_slide.shapes:
            _copy_shape_xml(shape, target_slide)
        
        # Заменить плейсхолдеры на реальный текст
        _replace_in_text_elements(target_slide._element, content)
        
        # Заменить картинку если есть
        if "_image_path" in content:
            _replace_image_in_slide(target_slide, content["_image_path"])
        
        # Очистить оставшиеся незаполненные плейсхолдеры
        _clear_missing_placeholders(target_slide._element)
    
    # Сохранить
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    output_prs.save(output_path)
    return output_path


if __name__ == "__main__":
    # Тест — собрать мини-презу из 3 слайдов разных типов
    test_spec = [
        {
            "layout_kind": "title",
            "content": {
                "DECK_LABEL": "ТЕСТ",
                "TITLE": "Заголовок",
                "SUBTITLE": "Это тест сборки",
                "YEAR": "2025",
            },
        },
        {
            "layout_kind": "bullets",
            "content": {
                "HEADING": "Ключевые пункты",
                "DECK_LABEL": "ТЕСТ",
                "BULLET_1_TITLE": "Первый",
                "BULLET_1_TEXT": "Описание первого пункта",
                "BULLET_2_TITLE": "Второй",
                "BULLET_2_TEXT": "Описание второго",
                "BULLET_3_TITLE": "Третий",
                "BULLET_3_TEXT": "Описание третьего",
            },
        },
        {
            "layout_kind": "closing",
            "content": {
                "CLOSING_HEADLINE": "Конец.",
                "CLOSING_SUBTITLE": "Вопросы?",
                "CONTACTS": "test@example.com",
            },
        },
    ]
    
    for style in ["formal", "corporate", "bold"]:
        output = f"/tmp/test_build_{style}.pptx"
        build_presentation(test_spec, style, output)
        print(f"Built: {output}")
