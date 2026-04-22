"""
Unit test для slide_builder — прогоняет все 11 layouts × 3 styles = 33 комбинации.
Не требует ANTHROPIC_API_KEY, тестирует только сборку.

Запуск: REDESIGNER_MASTERS_DIR=/path/to/masters python3 test_builder.py
"""
import os
import sys
import tempfile
from pathlib import Path

# Relative import support
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

try:
    from .slide_builder import build_presentation
    from .layout_catalog import LAYOUT_KINDS
except ImportError:
    from slide_builder import build_presentation
    from layout_catalog import LAYOUT_KINDS


# Тестовый контент для каждого layout'а
TEST_CONTENT = {
    "title": {
        "DECK_LABEL": "ТЕСТ", "YEAR": "2025",
        "TITLE": "Test Title",
        "TITLE_ACCENT": "accent",
        "SUBTITLE": "Тестовый подзаголовок",
        "SUPER_LABEL": "SUPER LABEL",
    },
    "section": {
        "DECK_LABEL": "ТЕСТ", "YEAR": "2025",
        "SECTION_NUM": "01",
        "SECTION_TITLE": "Раздел первый",
        "SECTION_DESCRIPTION": "Описание раздела для тестирования верстки.",
    },
    "text_heavy": {
        "DECK_LABEL": "ТЕСТ", "YEAR": "2025",
        "HEADING": "Большая тема",
        "LEAD": "Подзаголовок для контекста",
        "BODY_TEXT": "Основной текст абзаца с описанием ситуации. Достаточно длинный чтобы проверить переносы и плотность.",
    },
    "bullets": {
        "DECK_LABEL": "ТЕСТ", "YEAR": "2025",
        "HEADING": "5 пунктов",
        "BULLET_1_NUM": "01", "BULLET_1_TITLE": "Первый", "BULLET_1_TEXT": "Описание первого",
        "BULLET_2_NUM": "02", "BULLET_2_TITLE": "Второй", "BULLET_2_TEXT": "Описание второго",
        "BULLET_3_NUM": "03", "BULLET_3_TITLE": "Третий", "BULLET_3_TEXT": "Описание третьего",
        "BULLET_4_NUM": "04", "BULLET_4_TITLE": "Четвёртый", "BULLET_4_TEXT": "Описание",
        "BULLET_5_NUM": "05", "BULLET_5_TITLE": "Пятый", "BULLET_5_TEXT": "Описание",
    },
    "two_columns": {
        "DECK_LABEL": "ТЕСТ", "YEAR": "2025",
        "HEADING": "Сравнение",
        "LEFT_TITLE": "План", "LEFT_SUB": "Цели", "LEFT_TEXT": "Левая колонка — план на год.",
        "RIGHT_TITLE": "Факт", "RIGHT_SUB": "Итоги", "RIGHT_TEXT": "Правая колонка — факт.",
    },
    "stats": {
        "DECK_LABEL": "ТЕСТ", "YEAR": "2025",
        "HEADING": "Метрики",
        "SUBTITLE": "Ключевые числа",
        "STAT_1_CATEGORY": "ВЫРУЧКА", "STAT_1_VALUE": "2.8B", "STAT_1_LABEL": "млрд ₽",
        "STAT_2_CATEGORY": "РОСТ", "STAT_2_VALUE": "+18%", "STAT_2_LABEL": "к прошлому",
        "STAT_3_CATEGORY": "КЛИЕНТЫ", "STAT_3_VALUE": "247", "STAT_3_LABEL": "активных",
        "STAT_4_CATEGORY": "ШТАТ", "STAT_4_VALUE": "132", "STAT_4_LABEL": "в команде",
    },
    "quote": {
        "DECK_LABEL": "ТЕСТ", "YEAR": "2025",
        "QUOTE_TEXT": "Мы делаем то, что обещаем.",
        "QUOTE_AUTHOR": "Иван Петров",
        "QUOTE_ROLE": "Управляющий партнёр",
    },
    "table": {
        "DECK_LABEL": "ТЕСТ", "YEAR": "2025",
        "HEADING": "Сравнение",
        "COL_1_HEADER": "Показатель", "COL_2_HEADER": "План", "COL_3_HEADER": "Факт", "COL_4_HEADER": "%",
        "R1_C1": "Выручка", "R1_C2": "2.4B", "R1_C3": "2.8B", "R1_C4": "117%",
        "R2_C1": "Клиенты", "R2_C2": "200", "R2_C3": "247", "R2_C4": "124%",
        "R3_C1": "EBITDA", "R3_C2": "680M", "R3_C3": "742M", "R3_C4": "109%",
        "R4_C1": "Штат", "R4_C2": "125", "R4_C3": "132", "R4_C4": "106%",
    },
    "closing": {
        "CLOSING_HEADLINE": "Спасибо.",
        "CLOSING_SUBTITLE": "Готовы к диалогу",
        "CONTACTS": "test@company.ru",
    },
    "image_text": {
        "DECK_LABEL": "ТЕСТ",
        "KICKER": "КЛЮЧЕВОЕ",
        "HEADING": "Заголовок с картинкой",
        "BODY_TEXT": "Описание рядом с изображением.",
        # _image_path не передан — оставим placeholder из master'а
    },
    "image_full": {
        "DECK_LABEL": "ТЕСТ",
        "KICKER": "КЛЮЧЕВОЕ",
        "HEADING": "Full-bleed",
        "BODY_TEXT": "Текст поверх изображения.",
    },
}


def test_all_combinations():
    """Тест: собрать презу для каждого (style × layout) и проверить что файл создался."""
    styles = ["formal", "corporate", "bold"]
    layouts = LAYOUT_KINDS  # 11 лейаутов
    
    results = {"passed": 0, "failed": 0, "errors": []}
    
    for style in styles:
        for layout_kind in layouts:
            if layout_kind not in TEST_CONTENT:
                continue
            
            spec = [{
                "layout_kind": layout_kind,
                "content": TEST_CONTENT[layout_kind],
            }]
            
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                tmp_path = tmp.name
            
            try:
                build_presentation(spec, style, tmp_path)
                # Проверить что файл создан и размер > 1KB
                size = os.path.getsize(tmp_path)
                if size < 1024:
                    raise RuntimeError(f"File too small: {size} bytes")
                results["passed"] += 1
                os.unlink(tmp_path)
            except Exception as e:
                results["failed"] += 1
                results["errors"].append(f"{style}/{layout_kind}: {e}")
                if os.path.exists(tmp_path):
                    os.unlink(tmp_path)
    
    return results


def test_partial_bullets():
    """Тест: передать только 3 пункта из 5 — BULLET_4/5 должны быть скрыты."""
    spec = [{
        "layout_kind": "bullets",
        "content": {
            "HEADING": "Только 3",
            "BULLET_1_NUM": "01", "BULLET_1_TITLE": "A", "BULLET_1_TEXT": "a",
            "BULLET_2_NUM": "02", "BULLET_2_TITLE": "B", "BULLET_2_TEXT": "b",
            "BULLET_3_NUM": "03", "BULLET_3_TITLE": "C", "BULLET_3_TEXT": "c",
        },
    }]
    
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp_path = tmp.name
    
    try:
        build_presentation(spec, "formal", tmp_path)
        # Проверить что пустые плейсхолдеры очистились
        from pptx import Presentation
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        all_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += shape.text_frame.text + " "
        
        # В тексте не должно быть {BULLET_4_NUM} и т.п.
        assert "{BULLET_4_NUM}" not in all_text, "BULLET_4_NUM placeholder not cleared"
        assert "{BULLET_5_NUM}" not in all_text, "BULLET_5_NUM placeholder not cleared"
        assert "01" in all_text, "BULLET_1_NUM not filled"
        
        os.unlink(tmp_path)
        return True
    except Exception as e:
        print(f"FAIL partial_bullets: {e}")
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)
        return False


def test_multi_slide():
    """Тест: собрать многослайдовую презу (смешанные layouts)."""
    spec = [
        {"layout_kind": "title", "content": TEST_CONTENT["title"]},
        {"layout_kind": "section", "content": TEST_CONTENT["section"]},
        {"layout_kind": "bullets", "content": TEST_CONTENT["bullets"]},
        {"layout_kind": "stats", "content": TEST_CONTENT["stats"]},
        {"layout_kind": "closing", "content": TEST_CONTENT["closing"]},
    ]
    
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp_path = tmp.name
    
    try:
        build_presentation(spec, "corporate", tmp_path)
        from pptx import Presentation
        prs = Presentation(tmp_path)
        assert len(prs.slides) == 5, f"Expected 5 slides, got {len(prs.slides)}"
        os.unlink(tmp_path)
        return True
    except Exception as e:
        print(f"FAIL multi_slide: {e}")
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)
        return False


if __name__ == "__main__":
    print("=" * 60)
    print("Testing builder: all 3 styles × 11 layouts = 33 combinations")
    print("=" * 60)
    r = test_all_combinations()
    print(f"\nResults: {r['passed']} passed, {r['failed']} failed")
    if r["errors"]:
        print("\nErrors:")
        for err in r["errors"]:
            print(f"  - {err}")
    
    print("\n" + "=" * 60)
    print("Testing partial bullets (3 of 5)")
    print("=" * 60)
    ok = test_partial_bullets()
    print(f"Partial bullets: {'PASS' if ok else 'FAIL'}")
    
    print("\n" + "=" * 60)
    print("Testing multi-slide presentation")
    print("=" * 60)
    ok = test_multi_slide()
    print(f"Multi-slide: {'PASS' if ok else 'FAIL'}")
