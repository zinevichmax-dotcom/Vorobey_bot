"""
Парсер авторской PPTX — извлекает структурированный контент каждого слайда.

Для каждого слайда возвращает:
{
  "index": int,
  "text_blocks": [str, ...],         # список всех текстовых блоков в порядке появления
  "title_candidate": str,            # самый крупный/жирный текст как кандидат в заголовок
  "bullet_candidates": [str, ...],   # если есть список буллетов
  "images": [{"path": str, "width_emu": int, "height_emu": int}, ...],
  "tables": [{"rows": int, "cols": int, "data": [[str]]}],
  "has_chart": bool,
  "has_smartart": bool,
  "has_media": bool,                 # video/audio
  "total_chars": int,
  "raw_text": str,                   # весь текст склеенный
}
"""
import os
import shutil
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_image_from_shape(shape, output_dir: str, slide_idx: int, img_idx: int) -> dict | None:
    """Извлечь изображение из shape и сохранить в файл."""
    try:
        image = shape.image
        ext = image.ext  # 'png', 'jpg', etc.
        filename = f"slide{slide_idx}_img{img_idx}.{ext}"
        filepath = os.path.join(output_dir, filename)
        with open(filepath, "wb") as f:
            f.write(image.blob)
        return {
            "path": filepath,
            "width_emu": shape.width,
            "height_emu": shape.height,
            "left_emu": shape.left,
            "top_emu": shape.top,
        }
    except Exception as e:
        return None


def extract_table(shape) -> dict:
    """Извлечь таблицу."""
    table = shape.table
    rows_data = []
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            row_data.append(cell.text.strip())
        rows_data.append(row_data)
    return {
        "rows": len(rows_data),
        "cols": len(rows_data[0]) if rows_data else 0,
        "data": rows_data,
    }


def parse_pptx(pptx_path: str, images_output_dir: str) -> dict:
    """
    Распарсить PPTX. Возвращает структуру:
    {
      "total_slides": int,
      "slide_width_emu": int,
      "slide_height_emu": int,
      "slides": [ ... ]   # список slide_info
    }
    """
    os.makedirs(images_output_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    
    result = {
        "total_slides": len(prs.slides),
        "slide_width_emu": prs.slide_width,
        "slide_height_emu": prs.slide_height,
        "slides": [],
    }
    
    for slide_idx, slide in enumerate(prs.slides):
        text_blocks = []
        images = []
        tables = []
        has_chart = False
        has_smartart = False
        has_media = False
        img_counter = 0
        
        # Для определения title candidate — самый крупный текст
        largest_text = ""
        largest_size = 0
        
        for shape in slide.shapes:
            # Текст
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    text_blocks.append(text)
                    # Найти максимальный размер шрифта
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            sz = run.font.size.pt if run.font.size else 18
                            if sz > largest_size and len(text) < 200:
                                largest_size = sz
                                largest_text = text
            
            # Картинка
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_info = extract_image_from_shape(shape, images_output_dir, slide_idx, img_counter)
                if img_info:
                    images.append(img_info)
                    img_counter += 1
            
            # Таблица
            if shape.has_table:
                tables.append(extract_table(shape))
            
            # График
            if shape.has_chart:
                has_chart = True
            
            # SmartArt (проверка через графический fragment внутри shape'а)
            try:
                xml_elem = shape.element
                # SmartArt использует dgm: namespace или graphicFrame с diagramId
                xml_str = etree.tostring(xml_elem, encoding="unicode")
                if ('xmlns:dgm=' in xml_str or 
                    'dgm:relIds' in xml_str or 
                    '<dgm:' in xml_str):
                    has_smartart = True
            except Exception:
                pass
            
            # Медиа (видео/аудио)
            if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                has_media = True
        
        # Попробуем определить буллет-список
        bullet_candidates = []
        for block in text_blocks:
            # Текстовый блок с несколькими строками, где каждая короткая — это буллеты
            lines = [l.strip() for l in block.split("\n") if l.strip()]
            if len(lines) >= 2 and all(len(l) < 200 for l in lines):
                bullet_candidates = lines
                break
        
        raw_text = "\n".join(text_blocks)
        total_chars = len(raw_text)
        
        result["slides"].append({
            "index": slide_idx,
            "text_blocks": text_blocks,
            "title_candidate": largest_text if largest_size >= 24 else (text_blocks[0] if text_blocks else ""),
            "bullet_candidates": bullet_candidates,
            "images": images,
            "tables": tables,
            "has_chart": has_chart,
            "has_smartart": has_smartart,
            "has_media": has_media,
            "total_chars": total_chars,
            "raw_text": raw_text,
        })
    
    return result


if __name__ == "__main__":
    # Тест парсера на Muted оригинале
    test_file = "/home/claude/refs/4_muted.pptx"
    output = "/tmp/parsed_images"
    if os.path.exists(test_file):
        data = parse_pptx(test_file, output)
        print(f"Total slides: {data['total_slides']}")
        print(f"Size: {Emu(data['slide_width_emu']).inches:.1f}x{Emu(data['slide_height_emu']).inches:.1f} inches")
        for slide in data['slides'][:3]:
            print(f"\nSlide {slide['index']}:")
            print(f"  title_candidate: {slide['title_candidate'][:60]}")
            print(f"  text_blocks: {len(slide['text_blocks'])}")
            print(f"  images: {len(slide['images'])}")
            print(f"  tables: {len(slide['tables'])}")
            print(f"  has_chart: {slide['has_chart']}")
            print(f"  has_smartart: {slide['has_smartart']}")
            print(f"  total_chars: {slide['total_chars']}")
