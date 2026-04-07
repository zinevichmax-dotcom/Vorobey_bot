"""
Парсер PPTX: извлекает структуру слайдов в JSON.
Не меняет файл — только читает.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
import json


def parse_pptx(file_path: str) -> dict:
    """
    Извлекает полную структуру презентации.
    Возвращает dict с метаданными + слайды + элементы.
    """
    prs = Presentation(file_path)
    slide_width = prs.slide_width  # EMU
    slide_height = prs.slide_height

    slides = []
    for i, slide in enumerate(prs.slides):
        shapes_data = []
        for shape in slide.shapes:
            shape_info = {
                "shape_id": shape.shape_id,
                "name": shape.name,
                "left": shape.left,     # EMU
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
                "type": _get_shape_type(shape),
            }

            # Текстовые данные
            if shape.has_text_frame:
                paragraphs = []
                for para in shape.text_frame.paragraphs:
                    runs = []
                    for run in para.runs:
                        runs.append({
                            "text": run.text,
                            "font_name": run.font.name,
                            "font_size": run.font.size,  # EMU или None
                            "bold": run.font.bold,
                            "italic": run.font.italic,
                            "color": _get_color(run.font.color),
                        })
                    paragraphs.append({
                        "text": para.text,
                        "alignment": str(para.alignment),
                        "level": para.level,
                        "runs": runs,
                    })
                shape_info["paragraphs"] = paragraphs

            # Таблица
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = []
                    for cell in row.cells:
                        cells.append(cell.text)
                    rows.append(cells)
                shape_info["table"] = {
                    "rows": len(table.rows),
                    "cols": len(table.columns),
                    "data": rows,
                }

            # Изображение
            if hasattr(shape, "image"):
                shape_info["has_image"] = True

            shapes_data.append(shape_info)

        slides.append({
            "index": i,
            "shapes": shapes_data,
        })

    return {
        "slide_count": len(prs.slides),
        "slide_width": slide_width,
        "slide_height": slide_height,
        "slides": slides,
    }


def _get_shape_type(shape) -> str:
    """Определяет тип шейпа."""
    if hasattr(shape, "image"):
        return "image"
    if shape.has_table:
        return "table"
    if shape.has_text_frame:
        # Пытаемся понять: заголовок или тело
        if "title" in shape.name.lower():
            return "title"
        if "subtitle" in shape.name.lower():
            return "subtitle"
        return "text"
    return "other"


def _get_color(color_obj) -> str | None:
    """Безопасно извлекает цвет."""
    try:
        if color_obj and color_obj.rgb:
            return str(color_obj.rgb)
    except Exception:
        pass
    return None


# Тест
if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        result = parse_pptx(sys.argv[1])
        print(json.dumps(result, indent=2, default=str, ensure_ascii=False))
    else:
        print("Использование: python pptx_parser.py файл.pptx")